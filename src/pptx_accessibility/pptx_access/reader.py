"""PPTX Reader - wrapper around python-pptx for reading presentations."""

from pathlib import Path
from typing import Any

from loguru import logger
from pptx import Presentation
from pptx.util import Pt


class PPTXReader:
    """Reads PPTX presentations and extracts accessibility-relevant information."""

    def __init__(self, file_path: Path | str) -> None:
        """Initialize PPTX reader.

        Args:
            file_path: Path to PPTX file
        """
        self.file_path = Path(file_path)
        self.presentation = Presentation(str(self.file_path))
        logger.info(f"Loaded PPTX: {self.file_path.name} with {len(self.presentation.slides)} slides")

    def get_slide_count(self) -> int:
        """Get total number of slides."""
        return len(self.presentation.slides)

    def get_slide_title(self, slide_index: int) -> str | None:
        """Get title of a specific slide.

        Args:
            slide_index: 0-based slide index

        Returns:
            Slide title or None if no title found
        """
        slide = self.presentation.slides[slide_index]

        # Check if slide has a title placeholder
        if slide.shapes.title:
            return slide.shapes.title.text.strip() if slide.shapes.title.text else None

        return None

    def get_all_text_shapes(self, slide_index: int) -> list[dict[str, Any]]:
        """Get all text shapes from a slide.

        Args:
            slide_index: 0-based slide index

        Returns:
            List of dictionaries with shape information
        """
        slide = self.presentation.slides[slide_index]
        text_shapes = []

        for shape_idx, shape in enumerate(slide.shapes):
            if not hasattr(shape, "text"):
                continue

            text = shape.text.strip()
            if not text:
                continue

            shape_info = {
                "shape_index": shape_idx,
                "shape_id": shape.shape_id,
                "shape_name": shape.name,
                "text": text,
                "is_title": shape == slide.shapes.title,
                "is_placeholder": shape.is_placeholder if hasattr(shape, "is_placeholder") else False,
                "text_frame": self._extract_text_frame_info(shape) if hasattr(shape, "text_frame") else None,
            }

            text_shapes.append(shape_info)

        return text_shapes

    def _extract_text_frame_info(self, shape: Any) -> dict[str, Any]:
        """Extract text formatting information from text frame.

        Args:
            shape: Shape with text_frame

        Returns:
            Dictionary with text frame information
        """
        if not hasattr(shape, "text_frame"):
            return {}

        text_frame = shape.text_frame
        paragraphs_info = []

        for para in text_frame.paragraphs:
            for run in para.runs:
                run_info = {
                    "text": run.text,
                    "font_name": run.font.name if run.font.name else "Unknown",
                    "font_size": run.font.size.pt if run.font.size else None,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "underline": run.font.underline,
                }

                # Try to get font color
                if run.font.color and run.font.color.rgb:
                    rgb = run.font.color.rgb
                    run_info["font_color"] = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                else:
                    run_info["font_color"] = None

                paragraphs_info.append(run_info)

        return {
            "paragraph_count": len(text_frame.paragraphs),
            "runs": paragraphs_info,
        }

    def get_all_images(self, slide_index: int) -> list[dict[str, Any]]:
        """Get all images from a slide.

        Args:
            slide_index: 0-based slide index

        Returns:
            List of image information dictionaries
        """
        slide = self.presentation.slides[slide_index]
        images = []

        for shape_idx, shape in enumerate(slide.shapes):
            # Check if shape is a picture
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_info = {
                    "shape_index": shape_idx,
                    "shape_id": shape.shape_id,
                    "shape_name": shape.name,
                    "has_alt_text": bool(shape.name and shape.name != "Picture"),
                    "alt_text": getattr(shape, "_element", None) and
                               getattr(shape._element, "desc", None) or None,
                }

                # Try to get image blob
                try:
                    if hasattr(shape, "image"):
                        image_info["content_type"] = shape.image.content_type
                        image_info["ext"] = shape.image.ext
                except Exception as e:
                    logger.debug(f"Could not extract image details: {e}")

                images.append(image_info)

        return images

    def has_title(self, slide_index: int) -> bool:
        """Check if slide has a title.

        Args:
            slide_index: 0-based slide index

        Returns:
            True if slide has a non-empty title
        """
        title = self.get_slide_title(slide_index)
        return title is not None and len(title) > 0

    def get_slide_layout_name(self, slide_index: int) -> str:
        """Get the layout name of a slide.

        Args:
            slide_index: 0-based slide index

        Returns:
            Layout name
        """
        slide = self.presentation.slides[slide_index]
        return slide.slide_layout.name

    def get_presentation_metadata(self) -> dict[str, Any]:
        """Get presentation metadata.

        Returns:
            Dictionary with metadata
        """
        core_props = self.presentation.core_properties

        return {
            "title": core_props.title or "Untitled",
            "author": core_props.author or "Unknown",
            "subject": core_props.subject,
            "keywords": core_props.keywords,
            "created": core_props.created,
            "modified": core_props.modified,
            "slide_count": self.get_slide_count(),
        }
